import os
import collections.abc
from dotenv import load_dotenv
from openai import OpenAI
import requests
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt

# Load environment variables
load_dotenv()

# Initialize OpenAI client
client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))

def generate_image(prompt):
    response = client.images.generate(
        model="dall-e-3",
        prompt=prompt,
        size="1024x1024",
        quality="standard",
        n=1,
    )
    image_url = response.data[0].url
    return image_url

def download_image(url):
    response = requests.get(url)
    return BytesIO(response.content)

def generate_image_prompt(slide_title, slide_content, presentation_summary):
    response = client.chat.completions.create(
        model="gpt-4",
        messages=[
            {"role": "system", "content": "You are an AI assistant that generates image prompts for presentation slides."},
            {"role": "user", "content": f"Generate a detailed image prompt for a slide with the following details:\n\nPresentation Summary: {presentation_summary}\n\nSlide Title: {slide_title}\n\nSlide Content: {slide_content}\n\nThe image should be relevant to the slide's content and the overall presentation theme."}
        ]
    )
    return response.choices[0].message.content

def create_slide(presentation, slide_title, slide_content, image_path, image_on_right, logo_path):
    slide_layout = presentation.slide_layouts[1]
    slide = presentation.slides.add_slide(slide_layout)

    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = slide_title
    
    # Add multiple bullet points
    text_frame = content.text_frame
    text_frame.clear()
    for point in slide_content:
        p = text_frame.add_paragraph()
        p.text = point
        p.level = 0

    # Set title font properties
    title.text_frame.paragraphs[0].runs[0].font.name = 'Calibri'
    title.text_frame.paragraphs[0].runs[0].font.size = Pt(22)
    title.text_frame.paragraphs[0].runs[0].font.bold = True

    # Set content font properties
    for paragraph in content.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = 'Calibri'
            run.font.size = Pt(20)

    image_width = Inches(4.5)
    margin = Inches(1.0)

    if image_on_right:
        left_position = presentation.slide_width - image_width
        content.width = presentation.slide_width - image_width - margin
        title.width = presentation.slide_width - image_width - margin * 2
        title.left = margin
        content.left = Inches(0.5)
        logo_left = Inches(0.5)
    else:
        left_position = 0
        content.width = presentation.slide_width - image_width - margin
        title.width = presentation.slide_width - image_width - margin * 2
        title.left = image_width + margin
        content.left = Inches(5)
        logo_left = presentation.slide_width - Inches(2.5)

    title.top = Inches(1)
    content.top = Inches(2.2)

    # Add image
    slide.shapes.add_picture(image_path, left_position, 0, width=image_width, height=presentation.slide_height)

    # Add logo
    logo_width = Inches(2)
    logo_height = Inches(0.75)
    logo_top = presentation.slide_height - logo_height - Inches(0.5)
    slide.shapes.add_picture(logo_path, logo_left, logo_top, width=logo_width, height=logo_height)

def create_cover_slide(presentation, title, subtitle, logo_path):
    slide_layout = presentation.slide_layouts[0]  # Using the title slide layout
    slide = presentation.slides.add_slide(slide_layout)

    title_placeholder = slide.shapes.title
    subtitle_placeholder = slide.placeholders[1]

    title_placeholder.text = title
    subtitle_placeholder.text = subtitle

    # Set font properties
    title_placeholder.text_frame.paragraphs[0].runs[0].font.name = 'Calibri'
    title_placeholder.text_frame.paragraphs[0].runs[0].font.size = Pt(44)
    title_placeholder.text_frame.paragraphs[0].runs[0].font.bold = True

    subtitle_placeholder.text_frame.paragraphs[0].runs[0].font.name = 'Calibri'
    subtitle_placeholder.text_frame.paragraphs[0].runs[0].font.size = Pt(32)

    # Add logo
    logo_width = Inches(3)
    logo_height = Inches(1.125)
    logo_left = (presentation.slide_width - logo_width) / 2
    logo_top = presentation.slide_height - logo_height - Inches(1)
    slide.shapes.add_picture(logo_path, logo_left, logo_top, width=logo_width, height=logo_height)

def main():
    presentation = Presentation()
    presentation.slide_width = Inches(10)
    presentation.slide_height = Inches(7.5)

    logo_path = 'logo.png'  # Make sure this file exists in the same directory as the script

    print("Welcome to the Interactive Presentation Generator!")
    
    # Create cover slide
    cover_title = input("Enter the presentation title for the cover slide: ")
    cover_subtitle = input("Enter the subtitle for the cover slide: ")
    create_cover_slide(presentation, cover_title, cover_subtitle, logo_path)

    slides_content = []
    image_on_right = True

    while True:
        choice = input("Press 1 to add a new slide or 2 to finish and create the presentation: ")
        
        if choice == '1':
            slide_title = input("Enter slide title: ")
            slide_content = []
            while True:
                content = input("Enter slide content (or press Enter to finish this slide): ")
                if content:
                    slide_content.append(content)
                else:
                    break
            slides_content.append((slide_title, slide_content))
        elif choice == '2':
            break
        else:
            print("Invalid choice. Please try again.")

    # Generate presentation summary
    presentation_summary = "\n".join([f"Slide {i+1}: {title}" for i, (title, _) in enumerate(slides_content)])
    summary_prompt = f"Summarize the following presentation outline:\n\n{presentation_summary}"
    
    summary_response = client.chat.completions.create(
        model="gpt-4",
        messages=[
            {"role": "system", "content": "You are an AI assistant that summarizes presentation outlines."},
            {"role": "user", "content": summary_prompt}
        ]
    )
    presentation_summary = summary_response.choices[0].message.content

    for i, (slide_title, slide_content) in enumerate(slides_content):
        print(f"Generating image for slide {i+1}...")
        image_prompt = generate_image_prompt(slide_title, "\n".join(slide_content), presentation_summary)
        image_url = generate_image(image_prompt)
        image_data = download_image(image_url)
        
        create_slide(presentation, slide_title, slide_content, image_data, image_on_right, logo_path)
        image_on_right = not image_on_right

    file_name = input("Enter the name for your presentation file (without extension): ")
    file_location = input("Enter the location to save your file: ")
    full_path = os.path.join(file_location, f"{file_name}.pptx")

    presentation.save(full_path)
    print(f"Presentation saved as {full_path}")

if __name__ == "__main__":
    main()